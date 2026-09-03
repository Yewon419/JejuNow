"""공식 기능설명서 양식(pptx)에 제주나우 내용을 채운다. 실행: python docs/contest/fill_pptx.py
산출: docs/contest/제주나우_기능설명서.pptx (PDF 변환은 export_pdf.ps1)"""

from __future__ import annotations

import copy
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.presentation import Presentation as PresentationDoc
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide
from pptx.table import Table, _Cell
from pptx.text.text import TextFrame
from pptx.util import Emu, Inches, Length, Pt

INK = RGBColor(0x1F, 0x29, 0x37)  # type: ignore[no-untyped-call]

HERE = Path(__file__).resolve().parent
TEMPLATE = HERE / "양식" / "2026 관광데이터 활용 공모전 웹앱 개발 부문 기능설명서 양식(작성용).pptx"
OUT = HERE / "제주나우_기능설명서.pptx"
CAP = HERE / "captures"
ICON = HERE.parent.parent / "app" / "public" / "icon-512.png"
PHONE_ASPECT = 1170 / 2532

TEAM_NAME = "(팀명: 콘텐츠랩 접수 팀명 그대로 입력)"

# ---------- 도우미 ----------


def set_text(
    tf: TextFrame,
    text: str,
    size: float,
    bold: bool = False,
    align: PP_PARAGRAPH_ALIGNMENT | None = None,
    line_spacing: float | None = None,
) -> None:
    tf.clear()  # type: ignore[no-untyped-call]
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()  # type: ignore[no-untyped-call]
        if align is not None:
            p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = INK  # 양식 셀 기본색이 흰색인 곳이 있어 명시


def set_cell(
    cell: _Cell,
    text: str,
    size: float = 12,
    bold: bool = False,
    align: PP_PARAGRAPH_ALIGNMENT | None = None,
) -> None:
    set_text(cell.text_frame, text, size, bold, align, line_spacing=1.1)


def remove_shape(shape: BaseShape) -> None:
    el = shape._element
    el.getparent().remove(el)


def shape_by_name(slide: Slide, name: str) -> GraphicFrame:
    for sh in slide.shapes:
        if sh.name == name and isinstance(sh, GraphicFrame):
            return sh
    raise KeyError(name)


def auto_shape_by_name(slide: Slide, name: str) -> BaseShape:
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(name)


def cell_box(shape: GraphicFrame, r: int, c: int, span_cols: int = 1) -> tuple[int, int, int, int]:
    tb = shape.table
    cols = [tb.columns[i].width for i in range(len(tb.columns))]
    rows = [tb.rows[i].height for i in range(len(tb.rows))]
    x = shape.left + sum(cols[:c])
    y = shape.top + sum(rows[:r])
    return x, y, sum(cols[c : c + span_cols]), rows[r]


def add_phone(slide: Slide, img: Path, box: tuple[int, int, int, int], height: Length) -> None:
    """세로 폰 캡처를 박스 안 가운데 위에 배치."""
    x, y, w, _h = box
    pic_w = int(height * PHONE_ASPECT)
    slide.shapes.add_picture(
        str(img), Emu(x + (w - pic_w) // 2), Emu(y + Inches(0.08)), height=height
    )


def duplicate_slide(prs: PresentationDoc, src: Slide) -> Slide:
    new: Slide = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):
        remove_shape(sh)
    for sh in src.shapes:
        new.shapes._spTree.insert_element_before(copy.deepcopy(sh._element), "p:extLst")
    return new


def keep_columns(table: Table, n: int, total_width: int) -> None:
    """앞 n개 열만 남기고 균등 폭으로. 남는 열은 gridCol·tc를 XML에서 제거."""
    grid = table._tbl.tblGrid
    for col in list(grid.gridCol_lst)[n:]:
        grid.remove(col)
    for tr in table._tbl.tr_lst:
        for tc in list(tr.tc_lst)[n:]:
            tr.remove(tc)
    for i in range(n):
        table.columns[i].width = Emu(total_width // n)


def delete_rows(table: Table, indices: list[int]) -> None:
    trs = list(table._tbl.tr_lst)
    for i in sorted(indices, reverse=True):
        table._tbl.remove(trs[i])


# ---------- 내용 ----------

SERVICE_NAME = "제주나우 (JEJU NOW)"
SERVICE_TYPE = (
    "웹 서비스 + 앱 서비스 (iOS)\n"
    "단일 코드베이스(Next.js + Capacitor)로 웹(https://jejunow.vercel.app)과 iOS 앱이 동일한 서비스"
)
SERVICE_SUMMARY = (
    "제주 관광지 804곳의 시간대별 혼잡도를 예측해 「지금 어디가 한적한가」를 알려주고, "
    "가려던 곳이 붐비면 같은 유형의 한적한 대안 관광지와 코스를 추천하는 분산 관광 서비스"
)
WHY = (
    "• 문제: 제주는 연 약 1,400만 명이 찾지만 수요가 성산일출봉·협재·우도 등 상위 10곳에 몰린다. "
    "특정 시간·장소 쏠림은 만족도 저하, 주차·도로 등 인프라 과부하, 탄소 배출로 이어진다.\n"
    "• 기존 앱의 한계: 인기·리뷰 기반 추천은 쏠림을 오히려 키운다. 여행자에게 필요한 정보는 "
    "「어디가 유명한가」가 아니라 「지금 어디가 한적한가」다.\n"
    "• 필요성: 제주특별자치도는 분산 관광을 정책 과제로 두고 있지만, 여행자가 현장에서 바로 쓸 수 "
    "있는 "
    "수요 분산 도구는 없다. 한국관광공사가 이미 보유한 관광지 정보(TourAPI)와 수요 통계(한국관광 "
    "데이터랩)를 "
    "여행자의 의사결정 순간에 연결하면, 추가 인프라 없이도 분산을 유도할 수 있다.\n"
    "• 정직한 설계: 실측 혼잡도가 없으므로 공개 통계를 학습한 예측값임을 앱 안에 명시하고, "
    "원천 데이터가 없는 관광지는 「추정」으로 표시한다."
)
FEATURES = (
    "1. 시간대별 혼잡도 예측 지도: 날짜·시간을 고르면 804곳 마커가 4단계(여유·보통·붐빔·혼잡) "
    "색으로 바뀐다. "
    "「한적한 곳만」 필터, 「추정치」 표시, 마커 탭 시 사진·시간대 그래프 시트\n"
    "2. 대안 관광지 추천: 붐빌 것으로 예측된 관광지와 같은 유형(TourAPI 소분류)의 한적한 곳을 "
    "가까운 순으로 제안, 한 번에 교체\n"
    "3. 여행 일정 혼잡도 시뮬레이션: 날짜·방문 순서를 넣으면 슬롯별 예측 혼잡도와 이동 "
    "거리·시간을 타임라인으로 표시, "
    "붐비는 슬롯에 시간 변경·대안 제안, 인앱 경로 지도(자동차·도보·대중교통)\n"
    "4. 오토플랜(자동 일정): 템포·이동수단·혼잡 허용도·여정 끝을 고르면 도달 가능성·유형 "
    "다양성·혼잡 회피를 만족하는 하루 일정을 자동 생성, "
    "2지선다 카드로 취향 학습\n"
    "5. 관광지 상세·경로 안내: 한국관광공사 OpenAPI 원천 정보(사진·소개·운영시간·전화·홈페이지)를 "
    "실시간 조회, "
    "시간대별 혼잡 예측 그래프와 「가장 한적한 시각」, 카카오맵·Apple 지도로 길찾기\n"
    "6. 홈(오늘의 제주): 내 일정 요약, 일정 스팟 근처 한적한 코스, 위치 권한 시 「내 주변 한적한 "
    "곳」 (위치는 기기 안에서만 사용)\n"
    "• 공통: 회원가입·로그인 없음, 일정은 기기 안에만 저장, 위치 권한을 거부해도 전 기능 사용 가능"
)
REGION = "제주특별자치도 (제주시·서귀포시 전역)"
REGION_IMPL = (
    "• 대상 관광지: TourAPI areaBasedList2를 제주 법정동 코드(lDongRegnCd=50)로 조회해 "
    "관광지·문화시설·레포츠 804곳 수집. "
    "운영시간 722곳, 공식 홈페이지 579곳 보강\n"
    "• 학습 데이터: 한국관광 데이터랩 인기관광지 통계를 제주시(50110)·서귀포시(50130) 기초지자체 "
    "TOP30 × 연령대 6구간 × 101개월(2018-01~2026-05)로 수집한 36,360행. "
    "제주 두 지자체의 수요 구조만 학습\n"
    "• 기상 피처: 기상청 ASOS 제주 관측 월 기온·강수 101개월\n"
    "• 공항 앵커: 오토플랜은 제주국제공항 출·도착을 기본 앵커로 두고 동선을 구성(카카오 로컬 API "
    "공항 전용 검색)\n"
    "• 도로 접근점: 산·해안 관광지 590곳에 최근접 주차장 좌표(route_lat/lng)를 별도 저장해 경로 "
    "계산이 실제 진입점으로 되게 함\n"
    "• 확장 설계: 향후 전국 확대를 전제로 지역코드·지자체 범위만 바꾸면 되도록 만들었고, 데이터와 "
    "앵커만 제주에 특화"
)


@dataclass(frozen=True)
class Feature:
    title: str
    desc: str
    steps: list[tuple[str, str]]


CORE: list[Feature] = [
    Feature(
        title="시간대별 혼잡도 예측 지도",
        desc="날짜·시간을 고르면 제주 관광지 804곳의 예측 혼잡도를 지도 위 4단계 색 마커로 "
        "보여주는 기능입니다. 시간에 따라 같은 장소의 혼잡도가 어떻게 바뀌는지 한눈에 비교할 수 "
        "있습니다.",
        steps=[
            (
                "f1-1-map13.png",
                "① 지도 탭에서 날짜와 시간(9~20시) 슬라이더를 고른다. 마커가 여유·보통·붐빔·혼잡 "
                "4색으로 표시된다.",
            ),
            (
                "f1-2-quiet-only.png",
                "② 「한적한 곳만」 필터를 켜면 여유·보통 관광지만 남는다. 「추정치 포함」으로 "
                "원천 데이터가 없는 곳의 표시 여부를 정한다.",
            ),
            (
                "f1-3-marker-sheet.png",
                "③ 마커를 누르면 사진·시간대별 그래프 시트가 올라온다. 그래프의 시간을 누르면 "
                "지도 전체가 그 시간 기준으로 바뀐다.",
            ),
            (
                "f1-4-map18.png",
                "④ 슬라이더를 18시로 옮기면 낮에 붐비던 곳이 여유로 바뀐다. 장소를 포기하지 않고 "
                "시간을 바꾸는 선택이 가능하다.",
            ),
        ],
    ),
    Feature(
        title="대안 관광지 추천",
        desc="붐빌 것으로 예측된 관광지와 같은 유형(TourAPI 소분류)의 한적한 관광지를 거리순으로 "
        "제안하고 한 번에 교체하는 기능입니다. 해수욕장에는 해수욕장을, 오름에는 오름을 추천해 "
        "동선과 여행 성격이 유지됩니다.",
        steps=[
            (
                "f3-1-schedule.png",
                "① 일정 슬롯이 붐빔·혼잡으로 예측되면 카드 아래에 「같은 카테고리의 한적한 대안」 "
                "3곳이 거리·혼잡도와 함께 뜬다(FastAPI /alternatives 실시간 조회, 실패 시 "
                "사전계산값 폴백).",
            ),
            (
                "f2-2-replaced.png",
                "② 「교체」를 누르면 그 슬롯이 대안 관광지로 바뀌고 이동 거리·시간이 다시 "
                "계산된다. 예: 함덕해수욕장(붐빔) → 세기알해변(여유).",
            ),
        ],
    ),
    Feature(
        title="여행 일정 혼잡도 시뮬레이션",
        desc="여행 날짜와 방문 순서를 넣으면 각 슬롯의 예측 혼잡도와 이동 거리·시간을 "
        "타임라인으로 보여주고, 실제 경로를 앱 안 지도로 확인하는 기능입니다.",
        steps=[
            (
                "f3-1-schedule.png",
                "① 일정 탭에서 날짜를 고르고 관광지와 시간을 추가한다. 슬롯마다 예측 혼잡도(4단계 "
                "배지·막대)가 표시된다(POST /simulate 라이브 추론, 8초 초과 시 사전계산값 폴백).",
            ),
            (
                "f3-2-route.png",
                "② 슬롯 사이 「경로 보기」를 누르면 카카오 길찾기 경로가 지도에 그려지고 "
                "거리·시간이 자동차/도보/대중교통 탭으로 표시된다. Apple 지도·카카오맵으로 바로 "
                "열 수 있다.",
            ),
            (
                "f2-2-replaced.png",
                "③ 시간을 바꾸거나 대안으로 교체하면 타임라인 전체가 다시 계산된다. 일정은 날짜별 "
                "시안으로 기기 안에 저장되고 링크로 공유할 수 있다.",
            ),
        ],
    ),
    Feature(
        title="오토플랜(자동 일정 생성)",
        desc="여행 템포·이동수단·혼잡 허용도·여정 끝을 고르면 하루 일정을 자동으로 짜 주는 "
        "기능입니다. 도달 가능성·운영시간·유형 다양성·혼잡 회피를 만족하는 조합을 찾고, 2지선다 "
        "카드로 취향을 학습합니다.",
        steps=[
            (
                "f4-1-autoplan-form.png",
                "① 일정 탭의 「자동으로 짜기」에서 템포(여유·보통·빽빽), 이동수단, 혼잡 허용도, "
                "여정 끝(시작점·공항·다른 장소), 날짜를 고른다.",
            ),
            (
                "f4-2-autoplan-choice.png",
                "② 후보 관광지를 혼잡도·취향·거리·다양성으로 점수화해 2지선다 카드를 보여준다. "
                "고를 때마다 취향을 학습하고, 조건을 만족하는 후보가 없으면 반경·혼잡 허용을 "
                "단계적으로 완화한다.",
            ),
            (
                "f3-1-schedule.png",
                "③ 완성된 일정은 일정 탭에 시안으로 저장되어 혼잡도 시뮬레이션·경로 보기로 "
                "이어진다. 이동수단이 도보·대중교통이면 경로 시간도 그 기준으로 표시된다.",
            ),
        ],
    ),
    Feature(
        title="관광지 상세 · 실시간 원천 정보 · 경로 안내",
        desc="관광지의 사진·소개·운영시간·전화·홈페이지를 한국관광공사 OpenAPI에서 실시간으로 "
        "받아 보여주고, 시간대별 혼잡 예측 그래프와 길찾기를 제공하는 기능입니다.",
        steps=[
            (
                "f5-1-detail.png",
                "① 관광지를 열면 TourAPI 원천 정보(대표 사진·운영시간·홈페이지·주소)가 표시된다. "
                "소개·전화·홈페이지는 화면을 열 때마다 detailCommon2를 실시간 호출해 "
                "갱신한다(/api/tour-detail).",
            ),
            (
                "f5-2-detail-chart.png",
                "② 오늘의 시간대별 혼잡 예측 그래프와 「가장 한적한 시각 / 가장 붐비는 시각」, "
                "「비슷한데 더 한적한 곳」을 함께 보여준다. 「지금」 마커로 현재 시각을 표시한다.",
            ),
            (
                "f3-2-route.png",
                "③ 「일정에 넣기」로 시뮬레이션에 추가하고, 경로 화면에서 Apple 지도·카카오맵으로 "
                "길찾기를 실행한다. 지도 앱에는 관광지 좌표만 전달되며 이용자 위치는 보내지 "
                "않는다.",
            ),
        ],
    ),
]

APIS: list[tuple[str, str]] = [
    (
        "한국관광공사 국문 관광정보 서비스(KorService2) · 지역기반 관광정보 조회 areaBasedList2",
        "제주(lDongRegnCd=50) 관광지·문화시설·레포츠 804곳의 명칭·좌표·분류(cat1~3)·대표 "
        "이미지·주소 수집. 지도 마커·대안 추천·일정의 관광지 목록 원천. 매주 수요일 자동 "
        "동기화(GitHub Actions)",
    ),
    (
        "〃 · 소개정보 조회 detailIntro2",
        "관광지 운영시간(개장·폐장) 수집. 상세 화면 운영시간 표시와 시간대 혼잡 프로파일의 "
        "개·폐장 경계에 사용. 매주 미확보 관광지 보강",
    ),
    (
        "〃 · 공통정보 조회 detailCommon2",
        "관광지 소개(overview)·전화·공식 홈페이지(예매 안내). 상세 화면을 열 때마다 실시간 호출해 "
        "최신 원천 값으로 갱신하고, 실패 시 주간 동기화 값으로 폴백",
    ),
    (
        "〃 · 이미지정보 조회 detailImage2",
        "대표 이미지가 없는 관광지의 첫 이미지 확보. 홈 코스 카드·지도 시트·상세 히어로 사진의 "
        "원천",
    ),
]
OTHERS: list[tuple[str, str]] = [
    (
        "한국관광 데이터랩 인기관광지 통계 (파일데이터)",
        "기초지자체(제주시·서귀포시) 인기관광지 TOP30 × 연령대 6구간 × 2018-01~2026-05 월별 "
        "36,360행. 예측 모델(LightGBM)의 학습 타겟(관광지×월 인기 점유율 %). 데이터랩은 OpenAPI가 "
        "없어 CSV로 수집",
    ),
    (
        "기상청 지상(ASOS) 월 자료 OpenAPI (공공데이터포털)",
        "제주 관측소 월평균 기온·강수량 101개월. 야외 관광 수요와 직결되는 기상 피처",
    ),
    (
        "카카오 API: 카카오맵 JavaScript SDK · 카카오모빌리티 길찾기 · 카카오 로컬(키워드 장소 "
        "검색)",
        "지도·마커·현위치 표시 / 관광지 간 자동차 경로·거리·시간(Vercel 함수 경유) / 오토플랜 "
        "공항·출발지 검색. 그 외 공휴일·연휴 피처(holidays 라이브러리)",
    ),
]
DIFF = (
    "• 추천 기준의 전환: 인기도·리뷰가 아니라 「예측 혼잡도(한적함)」로 추천한다. 인기 추천은 "
    "쏠림을 키우고, 한적함 추천은 수요를 분산시킨다.\n"
    "• 데이터: 한국관광 데이터랩 8년치 수요 시계열을 LightGBM으로 학습하고, 관광지 804곳 × 45일 × "
    "9~20시 예측을 매주 자동 갱신한다. "
    "test hold-out(2026-01~05) 기준 MAE 0.46(점유율 %p), MAPE 14.9%, 상위 30% 랭킹 일치 83.5%.\n"
    "• 시간 단위: 일별이 아닌 시간대별. 장소를 포기하는 대신 「한적한 시간에 가는」 선택지를 "
    "준다.\n"
    "• 대안의 성격 유지: 같은 유형(TourAPI 소분류) 안에서 가까운 순으로 제안해 동선과 여행 성격이 "
    "깨지지 않는다.\n"
    "• 완성도·안정성: 웹·iOS 동일 코드베이스, 라이브 추론 실패 시 사전계산값 폴백, GitHub Actions "
    "3종(예측 갱신·관광지 동기화·DB 유지)으로 무중단 운영. App Store 심사 대응 완료.\n"
    "• 정직성과 개인정보: 예측값임을 앱·스토어에 명시, 원천 없는 703곳은 「추정」 표시. 계정·서버 "
    "저장 없음, 위치는 기기 안에서만 사용(위치기반서비스사업자 비신고 대상)."
)
PLAN = (
    "• 단기(2026 하반기): 실측 혼잡 지표(주차장 점유·입장 카운트 등) 확보 → 예측→실측 "
    "캘리브레이션. "
    "카테고리·요일 휴리스틱인 일중 프로파일을 실측 곡선으로 교체. 소셜 로그인(v1.1 준비 완료)으로 "
    "일정 동기화.\n"
    "• 중기(2027): 전국 확대. 코드·디자인이 지역 비의존이라 TourAPI 지역코드와 데이터랩 지자체 "
    "범위만 바꾸면 강원·부산·전주 등으로 확장. "
    "지역 관광공사·렌터카·숙박 플랫폼과 연계해 분산 효과를 정량화.\n"
    "• 장기: 관광 수요 예측 API를 지자체·여행사에 B2B로 제공, 지속가능 관광 지표 대시보드로 정책 "
    "활용. "
    "실시간 혼잡 데이터가 연동되면 예측·실측 하이브리드 모델로 전환."
)

# ---------- 채우기 ----------


def main() -> None:
    prs = Presentation(str(TEMPLATE))
    s1, s2, s3, s4, s5, s6, s7, s8, s9 = list(prs.slides)

    # 1 표지
    t = shape_by_name(s1, "표 5").table
    set_cell(t.cell(0, 1), TEAM_NAME, 16)
    set_cell(t.cell(1, 1), SERVICE_NAME, 16)

    # 2 서비스 소개
    t = shape_by_name(s2, "표 16").table
    set_cell(t.cell(0, 1), SERVICE_NAME, 14, bold=True)
    set_cell(t.cell(1, 1), SERVICE_TYPE, 12)
    set_cell(t.cell(2, 1), SERVICE_SUMMARY, 13)
    set_cell(t.cell(3, 1), WHY, 11)

    # 3 핵심기능
    t = shape_by_name(s3, "표 16").table
    set_cell(t.cell(0, 1), FEATURES, 11.5)

    # 4 이미지
    sh = shape_by_name(s4, "표 16")
    t = sh.table
    set_cell(t.cell(0, 1), "", 12)
    set_cell(t.cell(1, 1), "", 12)
    x, y, w, _h = cell_box(sh, 0, 1)
    s4.shapes.add_picture(
        str(ICON), Emu(x + Inches(0.25)), Emu(y + Inches(0.16)), height=Inches(1.7)
    )
    tb = s4.shapes.add_textbox(
        Emu(x + Inches(2.2)), Emu(y + Inches(0.35)), Inches(7.5), Inches(1.4)
    )
    set_text(
        tb.text_frame,
        "제주나우 앱 아이콘\n\n파란 위치 핀 안의 막대그래프: 「장소마다 시간대별 혼잡도가 "
        "다르다」는 서비스의 핵심을 나타냅니다.\n웹(PWA)·iOS 앱 공통 아이콘.",
        12,
    )
    x, y, w, _h = cell_box(sh, 1, 1)
    shots = [
        ("1-home.png", "홈 · 오늘의 제주"),
        ("f1-1-map13.png", "지도 · 시간대별 혼잡도"),
        ("f3-1-schedule.png", "일정 · 시뮬레이션·대안"),
        ("f4-2-autoplan-choice.png", "오토플랜 · 2지선다"),
        ("f5-1-detail.png", "상세 · 원천 정보·경로"),
    ]
    pic_h = Inches(2.95)
    pic_w = int(pic_h * PHONE_ASPECT)
    gap = Inches(0.35)
    sx = x + (w - (pic_w * 5 + gap * 4)) // 2
    for i, (f, cap) in enumerate(shots):
        px = sx + i * (pic_w + gap)
        s4.shapes.add_picture(str(CAP / f), Emu(px), Emu(y + Inches(0.1)), height=pic_h)
        tbx = s4.shapes.add_textbox(
            Emu(px - Inches(0.3)),
            Emu(y + Inches(0.1) + pic_h + Inches(0.02)),
            Emu(pic_w + Inches(0.6)),
            Inches(0.25),
        )
        set_text(tbx.text_frame, cap, 9, align=PP_ALIGN.CENTER)

    # 5 지역 특화
    remove_shape(auto_shape_by_name(s5, "직사각형 4"))
    t = shape_by_name(s5, "표 2").table
    set_cell(t.cell(0, 1), REGION, 14, bold=True)
    set_cell(t.cell(1, 1), REGION_IMPL, 11)

    # 6 핵심 기능 × 5 (원본 1 + 복제 4). 복제본은 끝에 붙으므로 6번 뒤로 옮긴다
    feature_slides = [s6] + [duplicate_slide(prs, s6) for _ in range(len(CORE) - 1)]
    lst = prs.slides._sldIdLst
    appended = list(lst)[9:]
    for k, el in enumerate(appended):
        lst.remove(el)
        lst.insert(6 + k, el)

    for idx, (slide, feat) in enumerate(zip(feature_slides, CORE), 1):  # noqa: B905
        remove_shape(auto_shape_by_name(slide, "직사각형 18"))
        head = shape_by_name(slide, "표 14").table
        set_cell(head.cell(0, 0), f"핵심 기능{idx}", 11, bold=True, align=PP_ALIGN.CENTER)
        set_cell(head.cell(0, 1), feat.title, 11, bold=True)
        set_cell(head.cell(1, 1), feat.desc, 10)
        flow = shape_by_name(slide, "표 7")
        tb = flow.table
        n = len(feat.steps)
        total_w = sum(tb.columns[i].width for i in range(len(tb.columns)))
        if n < 4:
            # 머리행(기능 흐름도)은 4열 병합 셀 → 남길 열 수에 맞춰 gridSpan 갱신
            keep_columns(tb, n, total_w)
            head_tc = tb.cell(0, 0)._tc
            head_tc.set("gridSpan", str(n))
        for c in range(n):
            set_cell(tb.cell(1, c), "", 10)
            set_cell(tb.cell(2, c), "", 10)
        for c, (img, text) in enumerate(feat.steps):
            p = CAP / img
            if p.exists():
                add_phone(slide, p, cell_box(flow, 1, c), Inches(2.85))
            set_cell(tb.cell(2, c), text, 9.5)

    # 7 OpenAPI
    t = shape_by_name(s7, "표 6").table
    for i, (name, desc) in enumerate(APIS):
        set_cell(t.cell(2 * i, 2), name, 11, bold=True)
        set_cell(t.cell(2 * i + 1, 2), desc, 10)
        t.rows[2 * i + 1].height = Inches(0.62)
    delete_rows(t, [8, 9])

    # 8 기타 데이터
    remove_shape(auto_shape_by_name(s8, "직사각형 5"))
    t = shape_by_name(s8, "표 6").table
    for i, (name, desc) in enumerate(OTHERS):
        set_cell(t.cell(2 * i, 2), name, 11, bold=True)
        set_cell(t.cell(2 * i + 1, 2), desc, 10)
        t.rows[2 * i + 1].height = Inches(0.62)

    # 9 차별성·발전계획
    t = shape_by_name(s9, "표 2").table
    set_cell(t.cell(0, 1), DIFF, 10.5)
    set_cell(t.cell(1, 1), PLAN, 10.5)

    prs.save(str(OUT))
    print("saved", OUT, "slides", len(prs.slides))


if __name__ == "__main__":
    main()
